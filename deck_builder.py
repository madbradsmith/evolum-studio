#!/usr/bin/env python3
#!V_1 UPDATED FOR NEW CLEAN APP SPACE

"""
deck_builder.py

Purpose:
- Full replacement deck builder
- Removes any full-slide overlay layer that can block or dull images
- Keeps images full-bleed and visible
- Preserves localized text boxes only
- Supports legacy CLI usage from the pipeline:
    python3 deck_builder.py /home/madbrad/app/slide_plan.json
- Also supports:
    python3 deck_builder.py --project /home/madbrad/app

New in this version:
- Reads approved_brain_output.json image_plan when available
- Uses brain-directed image queries / tags for smarter stock selection
- Keeps poster / user-uploaded image priority
- Falls back to legacy folder logic when no contextual stock match exists
- Updated for renamed stock folders
"""

from __future__ import annotations

import argparse
import concurrent.futures
import io
import json
import os
import re
import shutil
import tempfile
import threading
import urllib.request
from pathlib import Path
from typing import Optional

from PIL import Image, ImageFilter, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BASE_DIR = Path(__file__).resolve().parent
APP_DIR = Path(__file__).resolve().parent
PROJECT_DIR = Path(__file__).resolve().parent

_DAI_UID = os.environ.get("DAI_USER_ID", "")
_fal_image_count = 0
_fal_count_lock = threading.Lock()
_DAI_WORK_DIR = os.environ.get("DAI_WORK_DIR", "")


def _write_deck_builder_tokens() -> None:
    _f = Path(_DAI_WORK_DIR) / "pipeline_tokens.json" if _DAI_WORK_DIR else BASE_DIR / "pipeline_tokens.json"
    try:
        existing = json.loads(_f.read_text(encoding="utf-8")) if _f.exists() else {}
    except Exception:
        existing = {}
    existing["fal_images"] = _fal_image_count
    try:
        _f.write_text(json.dumps(existing), encoding="utf-8")
    except Exception:
        pass


_work_ctx = Path(_DAI_WORK_DIR) / "user_upload_context.json" if _DAI_WORK_DIR else None
if _work_ctx and _work_ctx.exists():
    UPLOAD_CONTEXT_PATH = str(_work_ctx)
else:
    _ctx_name = f"user_upload_context_{_DAI_UID}.json" if _DAI_UID else "user_upload_context.json"
    UPLOAD_CONTEXT_PATH = str(BASE_DIR / _ctx_name)
    if not Path(UPLOAD_CONTEXT_PATH).exists():
        UPLOAD_CONTEXT_PATH = str(BASE_DIR / "user_upload_context.json")


def load_user_context():
    try:
        with open(UPLOAD_CONTEXT_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    except:
        return {}

user_context = load_user_context()
user_poster = user_context.get("poster_filename", "")
USER_VISUAL_STYLE = (user_context.get("visual_style") or "live_action").strip().lower()
visuals_root = APP_DIR / "visuals"

_ctx_visuals_str = user_context.get("visuals_root", "")
USER_VISUALS_UPLOAD_ROOT = Path(_ctx_visuals_str) if _ctx_visuals_str else None

if USER_VISUALS_UPLOAD_ROOT:
    POSTER_PATH = str(USER_VISUALS_UPLOAD_ROOT / "poster" / user_poster) if user_poster else None
else:
    _uid_for_vis = _DAI_UID or "anon"
    POSTER_PATH = (
        str(APP_DIR / "visuals" / "user_uploaded" / _uid_for_vis / "poster" / user_poster)
        if user_poster else None
    )

DEFAULT_SLIDE_PLAN = APP_DIR / "slide_plan.json"
DEFAULT_BRAIN_OUTPUT = APP_DIR / "approved_brain_output.json"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TOP_RULE_Y = Inches(0.36)
TOP_RULE_H = Inches(0.05)

LAYOUT_THEMES = {
    # Warm near-black ink base + gold accent — EVOLUM brand palette
    "cinematic_grounded":      {"base": (10,9,7),   "base2": (20,17,14), "glow": (201,169,97,36),  "accent": (201,169,97),  "font": "Georgia"},
    "cinematic_high_tension":  {"base": (14,10,10), "base2": (28,18,14), "glow": (220,60,40,50),   "accent": (220,100,80),  "font": "Arial Black"},
    "contained_nocturnal":     {"base": (8,8,16),   "base2": (16,16,28), "glow": (60,80,200,40),   "accent": (100,140,220), "font": "Trebuchet MS"},
    "institutional_cinematic": {"base": (12,14,18), "base2": (24,28,34), "glow": (80,120,180,36),  "accent": (160,185,210), "font": "Verdana"},
    "storybook_satirical":     {"base": (20,16,10), "base2": (36,28,18), "glow": (220,180,80,44),  "accent": (220,180,80),  "font": "Georgia"},
    "neon_social_chaos":       {"base": (10,8,18),  "base2": (20,14,30), "glow": (180,40,240,50),  "accent": (180,80,240),  "font": "Trebuchet MS"},
    "athletic_prestige":       {"base": (8,14,20),  "base2": (16,26,36), "glow": (40,160,220,44),  "accent": (60,160,220),  "font": "Verdana"},
}

_active_theme: dict = LAYOUT_THEMES["cinematic_grounded"]


def rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


def _theme_font() -> str:
    return _active_theme.get("font", "Arial")


def clean(text) -> str:
    return " ".join(str(text or "").split()).strip()


def normalize_key(text: str) -> str:
    text = clean(text).lower()
    text = re.sub(r"\s*\(.*?\)\s*", " ", text)
    text = re.sub(r"[^a-z0-9]+", " ", text)
    return " ".join(text.split())


def tokenize(text: str) -> list[str]:
    return [t for t in re.findall(r"[a-z0-9]+", clean(text).lower()) if t]


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def next_output_path(output_dir: Path, label: str = "", uid: str = "") -> Path:
    uid_prefix = f"{uid}_" if uid else ""
    label_part = f"_{label}" if label else ""
    prefix = f"{uid_prefix}pitch_deck{label_part}"
    nums = []
    pattern = f"{prefix}_v*.pptx"
    for p in output_dir.glob(pattern):
        m = re.search(rf"{re.escape(prefix)}_v(\d+)\.pptx$", p.name)
        if m:
            nums.append(int(m.group(1)))
    nxt = max(nums) + 1 if nums else 1
    return output_dir / f"{prefix}_v{nxt}.pptx"


def resolve_paths(args) -> tuple[Path, Path, Path]:
    if args.project:
        project_dir = Path(args.project).expanduser().resolve()
        output_dir = project_dir / "output"
        output_dir.mkdir(parents=True, exist_ok=True)
        return project_dir / "slide_plan.json", project_dir / "visuals", output_dir

    slide_plan = Path(args.slide_plan).expanduser().resolve() if args.slide_plan else DEFAULT_SLIDE_PLAN.resolve()
    if _DAI_WORK_DIR:
        work_slide_plan = Path(_DAI_WORK_DIR) / "slide_plan.json"
        if work_slide_plan.exists():
            slide_plan = work_slide_plan
        # When running per-project, write the deck into the project's work dir
        # so api_build_deck can persist it as a project asset.
        per_project_output = Path(_DAI_WORK_DIR) / "visuals" / "output"
        per_project_output.mkdir(parents=True, exist_ok=True)
        return slide_plan, Path(_DAI_WORK_DIR) / "visuals", per_project_output
    output_dir = APP_DIR / "visuals" / "output"  # persistent disk mount
    output_dir.mkdir(parents=True, exist_ok=True)
    return slide_plan, APP_DIR / "visuals", output_dir


def load_brain_output(project_dir: Path) -> dict:
    if _DAI_WORK_DIR:
        work_candidate = Path(_DAI_WORK_DIR) / "approved_brain_output.json"
        if work_candidate.exists():
            return load_json(work_candidate)
    candidate = project_dir / "approved_brain_output.json"
    if candidate.exists():
        return load_json(candidate)
    if DEFAULT_BRAIN_OUTPUT.exists():
        return load_json(DEFAULT_BRAIN_OUTPUT)
    return {}


_image_usage_counts: dict = {}


def reset_image_selection_state() -> None:
    global _image_usage_counts
    _image_usage_counts = {}


def _image_usage_key(path: Path) -> str:
    try:
        return str(path.relative_to(APP_DIR / "visuals")).lower()
    except Exception:
        return str(path).lower()


def _image_use_count(path: Path) -> int:
    return _image_usage_counts.get(_image_usage_key(path), 0)


def _mark_image_used(path: Optional[Path]) -> None:
    if not path:
        return
    key = _image_usage_key(path)
    _image_usage_counts[key] = _image_usage_counts.get(key, 0) + 1


def _pick_candidate_with_repeat_control(candidates: list[Path], start_idx: int, last_used_name: str = "") -> tuple[Optional[Path], int]:
    if not candidates:
        return None, start_idx

    ordered = [candidates[(start_idx + offset) % len(candidates)] for offset in range(len(candidates))]

    # Pass 1: best case — brand new image, never adjacent repeat.
    for offset, candidate in enumerate(ordered):
        if candidate.name == last_used_name:
            continue
        if _image_use_count(candidate) == 0:
            return candidate, (start_idx + offset + 1) % len(candidates)

    # Pass 2: allow one reuse, but never more than twice in a deck.
    for offset, candidate in enumerate(ordered):
        if candidate.name == last_used_name:
            continue
        if _image_use_count(candidate) == 1:
            return candidate, (start_idx + offset + 1) % len(candidates)

    # Pass 3: emergency fallback — still never repeat adjacent slides,
    # and still block anything already used twice or more.
    for offset, candidate in enumerate(ordered):
        if candidate.name == last_used_name:
            continue
        if _image_use_count(candidate) < 2:
            return candidate, (start_idx + offset + 1) % len(candidates)

    # If the image pool is truly exhausted, return None so the caller can
    # fall through to a different selection source instead of repeating forever.
    return None, start_idx


def _top_visual_folder(path: Path) -> str:
    try:
        rel_parts = path.relative_to(APP_DIR / "visuals").parts
        return rel_parts[0] if rel_parts else ""
    except Exception:
        return path.parent.name


def _select_user_image(current_files: list[Path], lookup_key: str, last_used_name: str = "") -> Optional[Path]:
    if not current_files:
        return None, "none"

    candidates = sorted(current_files, key=lambda p: (
        _image_use_count(p),
        p.name.lower(),
    ))
    start_idx = _user_rotation_counters.get(lookup_key, 0)
    candidate, next_idx = _pick_candidate_with_repeat_control(candidates, start_idx, last_used_name)
    _user_rotation_counters[lookup_key] = next_idx
    return candidate


def _brain_image_instruction(brain_output: dict, slide_title: str, slide_number: int) -> Optional[dict]:
    image_plan = brain_output.get("image_plan", [])
    if not isinstance(image_plan, list):
        return None

    normalized_title = normalize_key(slide_title)

    for item in image_plan:
        if not isinstance(item, dict):
            continue
        if item.get("slide_number") == slide_number:
            return item

    for item in image_plan:
        if not isinstance(item, dict):
            continue
        if normalize_key(item.get("slide_title", "")) == normalized_title:
            return item

    if slide_number == 1:
        for item in image_plan:
            if not isinstance(item, dict):
                continue
            if normalize_key(item.get("slide_title", "")) == "title":
                return item

    return None


def _score_stock_file_against_tags(path: Path, tags: list[str]) -> int:
    try:
        rel_text = normalize_key(str(path.relative_to(APP_DIR / "visuals")))
    except Exception:
        rel_text = normalize_key(str(path))
    filename_text = normalize_key(path.stem)
    combined = f"{rel_text} {filename_text}"

    score = 0
    for tag in tags:
        tag_norm = normalize_key(tag)
        if not tag_norm:
            continue
        parts = tag_norm.split()

        if tag_norm in combined:
            score += 8

        matched_parts = sum(1 for part in parts if part in combined)
        score += matched_parts * 3

    return score


def resolve_image_options_for_slide(
    visuals_dir: Path,
    slide_info: dict,
    image_for_slide: Optional[Path],
    image_source: str,
    slide_title: str,
) -> list[dict]:
    exts = {".png", ".jpg", ".jpeg", ".webp", ".PNG", ".JPG", ".JPEG", ".WEBP"}
    raw_options = slide_info.get("image_options") or []
    resolved_options: list[dict] = []
    seen_paths: set[str] = set()

    def add_option(payload: dict) -> None:
        path_val = str(payload.get("image_path", "") or "").strip()
        if not path_val or path_val in seen_paths:
            return
        seen_paths.add(path_val)
        resolved_options.append(payload)

    if image_for_slide:
        add_option({
            "rank": 1,
            "option_id": "selected",
            "label": "Current Pick",
            "focus": "selected",
            "image_path": str(image_for_slide),
            "image_name": image_for_slide.name,
            "image_source": image_source,
        })

    if isinstance(raw_options, list):
        for option in raw_options:
            if not isinstance(option, dict):
                continue
            option_rank = int(option.get("rank", len(resolved_options) + 1) or (len(resolved_options) + 1))
            option_id = clean(option.get("option_id") or f"option_{option_rank}")
            option_label = clean(option.get("label") or f"Option {option_rank}")
            option_focus = clean(option.get("focus") or "alternate")
            option_url = str(option.get("image_url") or "").strip()

            # If the option already has a valid image file, carry it forward directly.
            # This preserves FAL-generated and previously-resolved stock options across rebuilds.
            existing_path_str = str(option.get("image_path") or "").strip()
            if existing_path_str:
                existing = Path(existing_path_str)
                if not existing.is_absolute():
                    existing = (APP_DIR / existing).resolve()
                if existing.exists():
                    add_option({
                        "rank": option_rank,
                        "option_id": option_id,
                        "label": option_label,
                        "focus": option_focus,
                        "image_path": str(existing),
                        "image_name": existing.name,
                        "image_source": str(option.get("image_source") or "preserved"),
                        "image_url": option_url,
                    })
                    if len(resolved_options) >= 4:
                        break
                    continue


FAL_API_KEY = ""  # FAL removed — visuals/ library is the only image source
EVOLUM_SESSION_ID = os.environ.get("EVOLUM_SESSION_ID", "shared")
TMDB_API_KEY = ""  # TMDB removed — no external comp films
TMDB_IMAGE_BASE = "https://image.tmdb.org/t/p/w500"

_user_rotation_counters: dict = {}

_SLIDE_VISUAL_CONCEPTS = {
    "logline":              "A short demo logline. Replace this sample slate with your own projects.",
    "synopsis":             "cinematic scene, atmospheric, narrative moment",
    "synopsis 2":           "cinematic scene, mid-shot, dramatic tension",
    "synopsis 3":           "cinematic scene, close-up, emotional intensity",
    "protagonist":          "cinematic portrait, single character, dramatic lighting, film still",
    "antagonist":           "cinematic portrait, menacing figure, dramatic shadows, film still",
    "supporting characters":"cinematic ensemble shot, multiple characters, film still",
    "world":                "cinematic landscape, establishing shot, rich environment",
    "hook":                 "cinematic close-up, tension, dramatic moment",
    "conflict":             "cinematic confrontation, dramatic tension, high stakes",
    "stakes":               "cinematic wide shot, weight of consequence, dramatic",
    "tone":                 "cinematic mood shot, atmospheric lighting, visual tone",
    "story engine":         "cinematic action, driving force, momentum",
    "reversal":             "cinematic turning point, dramatic shift, pivotal moment",
    "themes":               "cinematic symbolic imagery, thematic visual metaphor",
    "why this movie":       "cinematic wide shot, cultural moment, compelling imagery",
    "comparables":          "cinematic collage feel, prestige film aesthetic",
    "market projections":   "cinematic wide shot, commercial appeal, high production value",
    "closing statement":    "cinematic final frame, powerful, memorable",
    "closing":              "cinematic hero shot, wide epic frame, powerful final image, golden light",
}

_GENRE_STYLE = {
    "horror":       "dark, unsettling, atmospheric horror, shadows, practical effects aesthetic",
    "thriller":     "tense, noir-influenced, sharp contrast, suspenseful",
    "comedy":       "warm lighting, vibrant colors, playful composition",
    "drama":        "naturalistic lighting, intimate, emotionally grounded",
    "action":       "dynamic, kinetic energy, bold framing, high contrast",
    "sci-fi":       "futuristic, cool tones, technological, epic scale",
    "fantasy":      "magical, rich colors, otherworldly, painterly lighting",
    "romance":      "warm golden tones, soft focus, intimate, emotional",
    "documentary":  "gritty realism, candid, natural light, observational",
    "animation":    "stylized, vibrant, expressive, dynamic",
}


_PERIOD_KEYWORDS = [
    # (trigger keywords, content description, rendering style)
    (["medieval", "kingdom", "castle", "jester", "knight", "throne", "feudal", "valoria", "sorcery", "realm", "peasant", "dungeon", "jousting", "royal court", "nobleman"],
     "medieval kingdom setting, period-accurate costume, stone castle, torchlight",
     "oil painting style, medieval illuminated manuscript aesthetic, rich jewel tones, Renaissance painting"),
    (["space", "spaceship", "galaxy", "planet", "alien", "starship", "orbit"],
     "outer space, futuristic spacecraft, alien world",
     "sci-fi concept art, dramatic cosmic lighting, highly detailed illustration"),
    (["future", "cyberpunk", "neon", "dystopia", "android", "robot"],
     "near-future dystopian city, neon-lit streets, rain-slicked pavement",
     "cyberpunk digital art, neon illustration, cinematic concept art"),
    (["western", "frontier", "cowboy", "saloon", "sheriff"],
     "american frontier, dusty plains, period-accurate western setting",
     "painted western landscape, golden hour, cinematic wide oil painting"),
    (["victorian", "1800s", "19th century", "gaslight", "corset"],
     "victorian era street, period costume, gaslight atmosphere",
     "Victorian oil portrait style, dramatic chiaroscuro, 1800s painting aesthetic"),
    (["1920", "1930", "prohibition", "jazz age", "noir", "gangster"],
     "1930s city, prohibition era, art deco interior",
     "1930s noir illustration, art deco style, high contrast dramatic painting"),
    (["war", "wwii", "battlefield", "soldier", "trench"],
     "wartime battlefield, period military uniform, gritty realism",
     "wartime painted illustration, dramatic military art, painterly realism"),
    (["ancient", "roman", "roman empire", "greek", "egypt", "pyramid", "colosseum", "gladiator", "gladiatorial", "legion", "centurion", "toga", "senate", "chariot"],
     "ancient Roman setting, Colosseum arena, marble columns, Roman legionary armor, classical Roman architecture, not medieval",
     "neoclassical oil painting, ancient Roman epic art in the style of Jean-Léon Gérôme, classical antiquity scene, no Gothic architecture, no medieval castles"),
    (["pirate", "ship", "ocean", "sail", "treasure", "sea"],
     "tall ship, age of sail, ocean adventure, period-accurate nautical",
     "swashbuckling painted adventure, maritime oil painting, dramatic seascape"),
    (["animation", "animated", "cartoon", "pixar", "anime"],
     "stylized animated world, expressive characters",
     "vibrant stylized illustration, animated film concept art"),
]


def _detect_period_style(brain_output: dict):
    fields = [
        brain_output.get("world", ""),
        brain_output.get("genre", ""),
        brain_output.get("tone", ""),
        brain_output.get("logline", ""),
        brain_output.get("title", ""),
        brain_output.get("story_engine", ""),
        brain_output.get("tagline", ""),
        brain_output.get("synopsis", ""),
        brain_output.get("setting", ""),
    ]
    combined = " ".join(str(f) for f in fields if f).lower()
    for keywords, content_style, render_style in _PERIOD_KEYWORDS:
        if any(re.search(rf"\b{re.escape(k)}\b", combined) for k in keywords):
            return content_style, render_style
    return "", ""


_VISUAL_STYLE_PREFIX = {
    "animated": "2D animation style, cartoon illustration, vibrant colors, stylized characters, no photorealism, ",
    "illustrated": "concept art, painterly digital illustration, cinematic artwork, stylized realism, ",
    "live_action": "",
}


def _parse_demographics(description: str) -> str:
    """Extract race/age/build/gender from a character description for accurate FAL prompts."""
    desc = description.lower()
    parts = []

    for race in ["black", "african american", "white", "caucasian", "latina", "latino",
                 "hispanic", "asian", "middle eastern", "indigenous", "biracial"]:
        if re.search(rf'\b{re.escape(race)}\b', desc):
            parts.append(race)
            break

    age_m = re.search(
        r'\b(early|mid|late)?\s*(teens?|twenties|thirties|forties|fifties|sixties|seventies|\d{2}s?)\b',
        desc,
    )
    if age_m:
        parts.append(age_m.group(0).strip())

    for build in ["broad-shouldered", "broad shouldered", "muscular", "stocky",
                  "lean", "athletic", "heavyset", "heavy-set", "slender", "petite", "tall"]:
        clean_build = build.replace("-", " ")
        if clean_build in desc or build in desc:
            parts.append(clean_build)
            break

    if re.search(r'\b(he|his|him|man|guy|male|boy)\b', desc):
        parts.append("man")
    elif re.search(r'\b(she|her|hers|woman|girl|female|lady)\b', desc):
        parts.append("woman")

    return ", ".join(parts) if parts else ""

def build_image_prompt(slide_title: str, brain_output: dict, slide_body: str = "") -> str:
    normalized = normalize_key(slide_title)
    concept = _SLIDE_VISUAL_CONCEPTS.get(normalized, "cinematic scene, dramatic lighting")
    style_prefix = _VISUAL_STYLE_PREFIX.get(USER_VISUAL_STYLE, "")

    # Pull story-specific context from brain output for each slide type
    logline   = str(brain_output.get("logline", "") or "").strip()[:180]
    synopsis  = str(brain_output.get("synopsis", "") or "").strip()[:220]
    protagonist = str(brain_output.get("protagonist", "") or "").strip()
    protagonist_summary = str(brain_output.get("protagonist_summary", "") or "").strip()[:120]
    antagonist = str(brain_output.get("antagonist", "") or "").strip()
    antagonist_summary = str(brain_output.get("antagonist_summary", "") or "").strip()[:120]
    themes_raw = brain_output.get("themes", [])
    themes = ", ".join(themes_raw[:3]) if isinstance(themes_raw, list) else str(themes_raw or "")[:100]
    story_engine = str(brain_output.get("story_engine", "") or "").strip()[:120]

    # Prefer story-specific content over slide body for key slides
    story_scene = {
        "logline":    logline or slide_body,
        "synopsis":   synopsis or slide_body,
        "synopsis 2": synopsis or slide_body,
        "synopsis 3": synopsis or slide_body,
        "protagonist": f"{protagonist} — {protagonist_summary}" if protagonist else slide_body,
        "antagonist":  f"{antagonist} — {antagonist_summary}" if antagonist else slide_body,
        "themes":      themes or slide_body,
        "story engine": story_engine or slide_body,
        "hook":        logline or slide_body,
        "conflict":    synopsis or slide_body,
        "stakes":      story_engine or synopsis or slide_body,
    }.get(normalized, slide_body)

    _raw_scene = (story_scene or "").replace("\n", " ").strip()[:220]
    if len(_raw_scene) == 220 and " " in _raw_scene:
        _raw_scene = _raw_scene[:_raw_scene.rfind(" ")]
    scene = _raw_scene

    genre = str(brain_output.get("genre", "drama")).lower()
    genre_style = ""
    for g, style in _GENRE_STYLE.items():
        if g in genre:
            genre_style = style
            break
    if not genre_style:
        genre_style = "cinematic, naturalistic lighting"

    tone = str(brain_output.get("tone", "")).lower()
    world = str(brain_output.get("world", "")).replace("\n", " ").strip()

    period_content, period_render = _detect_period_style(brain_output)

    # Character portrait slides: inject demographics for accurate representation
    if normalized in {"protagonist", "antagonist"}:
        _char_sum = protagonist_summary if normalized == "protagonist" else antagonist_summary
        _char_name = protagonist if normalized == "protagonist" else antagonist
        _demographics = _parse_demographics(_char_sum) if _char_sum else ""
        _char_desc = _demographics if _demographics else (_char_name[:60] if _char_name else "single character")
        if period_render:
            prompt = (
                f"{style_prefix}{period_render}, cinematic portrait, {_char_desc}, "
                f"dramatic portrait lighting, {period_content}, "
                f"highly detailed, no text, no watermarks, 16:9 aspect ratio"
            )
        else:
            prompt = (
                f"{style_prefix}cinematic portrait, {_char_desc}, "
                f"dramatic portrait lighting, film still, {genre_style}, "
                f"professional cinematography, ultra-detailed, photorealistic, "
                f"no text, no watermarks, 16:9 aspect ratio"
            )
        print(f"🎨 FAL prompt [{slide_title}]: {prompt}")
        return prompt

    # Title slide gets a dedicated movie poster prompt
    film_title = str(brain_output.get("title", "")).strip()
    if film_title and normalize_key(slide_title) == normalize_key(film_title):
        protagonist = str(brain_output.get("protagonist", "")).strip()
        protagonist_summary = str(brain_output.get("protagonist_summary", "")).strip()
        _demographics = _parse_demographics(protagonist_summary) if protagonist_summary else ""
        char_hint = _demographics if _demographics else (protagonist[:60] if protagonist else "")
        tone_hint = f", {tone[:60]}" if tone else ""
        if period_render:
            prompt = (
                f"{style_prefix}{period_render}, {period_content}, "
                f"movie poster composition, lone hero figure, dramatic portrait lighting, "
                f"cinematic title card framing"
                f"{', ' + char_hint if char_hint else ''}"
                f"{tone_hint}, "
                f"highly detailed, no text, no watermarks, no logos, 16:9 aspect ratio"
            )
        else:
            prompt = (
                f"{style_prefix}movie poster composition, lone hero figure, dramatic portrait lighting, "
                f"cinematic title card framing"
                f"{', ' + char_hint if char_hint else ''}"
                f"{', ' + genre_style if genre_style else ''}"
                f"{tone_hint}, "
                f"professional film poster, ultra-detailed, no text, no watermarks, 16:9 aspect ratio"
            )
        print(f"🎨 FAL prompt [{slide_title}]: {prompt}")
        return prompt

    if period_render:
        scene_hint = f", scene depicting: {scene}" if scene else f", {concept}"
        tone_hint = f", {tone[:50]}" if tone else ""
        prompt = (
            f"{style_prefix}{period_render}, {period_content}{scene_hint}, "
            f"{genre_style}{tone_hint}, "
            f"highly detailed, dramatic lighting, no text, no watermarks, no logos, 16:9 aspect ratio"
        )
    else:
        scene_hint = f"scene depicting: {scene}, " if scene else f"{concept}, "
        visual_world = world if world and not world.startswith("feature /") else ""
        world_hint = f", {visual_world[:80]}" if visual_world else ""
        tone_hint = f", {tone[:60]}" if tone else ""
        photo_suffix = "professional film still, 35mm, shallow depth of field, ultra-detailed, photorealistic, " if not style_prefix else "ultra-detailed, "
        prompt = (
            f"{style_prefix}{scene_hint}{genre_style}{world_hint}{tone_hint}, "
            f"{photo_suffix}no text, no watermarks, 16:9 aspect ratio"
        )

    print(f"🎨 FAL prompt [{slide_title}]: {prompt}")
    return prompt


def generate_fal_image(prompt: str, cache_path: Path) -> Optional[Path]:
    global _fal_image_count
    if not FAL_API_KEY:
        return None
    if cache_path.exists():
        return cache_path

    import urllib.error
    url = "https://fal.run/fal-ai/flux/schnell"
    payload = json.dumps({
        "prompt": prompt,
        "image_size": "landscape_16_9",
        "num_inference_steps": 4,
        "num_images": 1,
        "enable_safety_checker": True,
    }).encode("utf-8")

    req = urllib.request.Request(
        url,
        data=payload,
        headers={
            "Authorization": f"Key {FAL_API_KEY}",
            "Content-Type": "application/json",
        },
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            result = json.loads(resp.read().decode("utf-8"))
        image_url = result["images"][0]["url"]
        cache_path.parent.mkdir(parents=True, exist_ok=True)
        urllib.request.urlretrieve(image_url, cache_path)
        with _fal_count_lock:
            _fal_image_count += 1
        print(f"✨ FAL generated image for prompt: {prompt[:60]}...")
        return cache_path
    except Exception as e:
        print(f"⚠️  FAL image generation failed: {e}")
        return None


def find_image_for_slide(
    visuals_dir: Path,
    deck_title: str,
    slide_title: str,
    slide_number: int,
    brain_output: Optional[dict] = None,
    last_used_name: str = "",
    slide_body: str = ""
) -> tuple[Optional[Path], str]:
    # When running inside a pipeline build use project-scoped uploads so images
    # from a different project never bleed in via the shared current/ folder.
    if _DAI_WORK_DIR:
        _wd = Path(_DAI_WORK_DIR)
        poster_dir = _wd / "poster"
        current_dir = _wd / "uploads"
    else:
        _upload_root = USER_VISUALS_UPLOAD_ROOT if USER_VISUALS_UPLOAD_ROOT else (visuals_dir / "user_uploaded" / (_DAI_UID or "anon"))
        poster_dir = _upload_root / "poster"
        current_dir = _upload_root / "current"
    exts = {".png", ".jpg", ".jpeg", ".webp", ".PNG", ".JPG", ".JPEG", ".WEBP"}

    poster_files = [p for p in poster_dir.glob("*") if p.suffix in exts] if poster_dir.exists() else []
    current_files = [p for p in current_dir.glob("*") if p.suffix in exts] if current_dir.exists() else []

    normalized_title = normalize_key(slide_title)
    image_instruction = _brain_image_instruction(brain_output or {}, slide_title, slide_number)

    non_title_slides = {
        "logline", "synopsis", "synopsis 2", "synopsis 3", "synopsis 4",
        "protagonist", "antagonist", "supporting characters", "world", "hook",
        "conflict", "conflict engine", "stakes", "tone", "story engine",
        "reversal", "theme", "themes", "why this movie", "why this film",
        "audience", "visual style", "comparables", "market position",
        "market projections", "director vision", "casting ideas", "production scope",
        "closing statement", "closing"
    }

    if normalized_title not in non_title_slides:
        for p in poster_files:
            print(f"🖼️ Using POSTER for '{deck_title}': {p}")
            return p, "poster"

    key_map = {
        "logline": ["logline", "frame_1", "opening"],
        "synopsis": ["synopsis", "frame_2", "scene"],
        "protagonist": ["protagonist", "character", "hero", "frame_3"],
        "world": ["world", "city", "location", "frame_4"],
        "hook": ["hook", "frame_5"],
        "conflict": ["conflict", "frame_6"],
        "stakes": ["stakes"],
        "tone": ["tone", "mood"],
        "story engine": ["engine"],
        "reversal": ["reversal", "turn"],
        "themes": ["themes", "theme"],
        "why this movie": ["why"],
    }

    lookup_key = "synopsis" if normalized_title.startswith("synopsis") else normalized_title

    if lookup_key in key_map:
        matched_files: list[Path] = []
        for needle in key_map[lookup_key]:
            for p in current_files:
                if needle.lower() in p.name.lower():
                    if p not in matched_files:
                        matched_files.append(p)

        if matched_files:
            selected = _select_user_image(matched_files, lookup_key, last_used_name)
            if selected:
                print(f"🖼️ Using USER image for '{slide_title}': {selected}")
                return selected, "user"

    if current_files:
        selected = _select_user_image(current_files, lookup_key, last_used_name)
        if selected:
            print(f"🖼️ Using USER image fallback for '{slide_title}': {selected}")
            return selected, "user_fallback"

    if FAL_API_KEY and brain_output:
        prompt = build_image_prompt(slide_title, brain_output, slide_body)
        cache_dir = visuals_dir.parent / "generated_images" / EVOLUM_SESSION_ID
        safe_title = re.sub(r"[^a-z0-9_]", "_", normalized_title)
        cache_path = cache_dir / f"{slide_number:02d}_{safe_title}.jpg"
        generated = generate_fal_image(prompt, cache_path)
        if generated:
            return generated, "fal_generated"



    return None, "none"


# Pre-computed (path, combined_search_text) pairs — keyed by id of the cached file list
_stock_precomputed_cache: dict[int, list[tuple]] = {}


def _score_combined(combined: str, normalized_tags: list[str]) -> int:
    score = 0
    for tag in normalized_tags:
        if not tag:
            continue
        if tag in combined:
            score += 8
        score += sum(3 for part in tag.split() if part in combined)
    return score


# Cache rendered base backgrounds — one expensive render per theme per process
_base_bg_cache: dict[int, bytes] = {}


def _render_base_bg(theme: dict) -> bytes:
    width_px, height_px = 640, 360
    b1 = theme["base"]
    b2 = theme["base2"]
    glow_color = theme["glow"]

    img = Image.new("RGB", (width_px, height_px), b1)
    draw = ImageDraw.Draw(img)
    for y in range(height_px):
        t = y / max(1, height_px - 1)
        r = int(b1[0] + (b2[0] - b1[0]) * t)
        g = int(b1[1] + (b2[1] - b1[1]) * t)
        b = int(b1[2] + (b2[2] - b1[2]) * t)
        draw.line((0, y, width_px, y), fill=(r, g, b))

    glow = Image.new("RGBA", (width_px, height_px), (0, 0, 0, 0))
    gdraw = ImageDraw.Draw(glow)
    gdraw.ellipse(
        (width_px * 0.18, height_px * 0.15, width_px * 0.82, height_px * 0.95),
        fill=glow_color,
    )
    glow = glow.filter(ImageFilter.GaussianBlur(40))
    img = Image.alpha_composite(img.convert("RGBA"), glow).convert("RGB")

    vignette = Image.new("L", (width_px, height_px), 0)
    vdraw = ImageDraw.Draw(vignette)
    vdraw.ellipse(
        (-width_px * 0.05, -height_px * 0.12, width_px * 1.05, height_px * 1.08),
        fill=220,
    )
    vignette = vignette.filter(ImageFilter.GaussianBlur(50))
    dark = Image.new("RGB", (width_px, height_px), (10, 10, 12))
    img = Image.composite(img, dark, vignette)

    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=80, optimize=True)
    return buf.getvalue()


def add_base_background(slide) -> None:
    theme = _active_theme
    cache_key = id(theme)
    if cache_key not in _base_bg_cache:
        _base_bg_cache[cache_key] = _render_base_bg(theme)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
    tmp.write(_base_bg_cache[cache_key])
    tmp.close()
    slide.shapes.add_picture(tmp.name, 0, 0, width=SLIDE_W, height=SLIDE_H)
    os.unlink(tmp.name)


def add_blur_background(slide, image_path: Optional[Path]) -> None:
    if not image_path or not image_path.exists():
        return

    with Image.open(image_path) as im:
        bg = im.convert("RGB")
        bg.thumbnail((320, 180))
        bg = bg.resize((640, 360))
        bg = bg.filter(ImageFilter.GaussianBlur(8))
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
        bg.save(tmp.name, format="JPEG", quality=65, optimize=True)

    slide.shapes.add_picture(str(tmp.name), 0, 0, width=SLIDE_W, height=SLIDE_H)
    import os; os.unlink(tmp.name)


def add_center_image(slide, image_path: Optional[Path], scale_factor: float = 0.68) -> None:
    if not image_path or not image_path.exists():
        return

    with Image.open(image_path) as im:
        img = im.convert("RGB")
        img.thumbnail((960, 540))
        img_w, img_h = img.size

        if img_w <= 0 or img_h <= 0:
            return

        scale = min(float(SLIDE_W) / img_w, float(SLIDE_H) / img_h) * scale_factor
        render_w = int(img_w * scale)
        render_h = int(img_h * scale)
        left = int((float(SLIDE_W) - render_w) / 2)
        top = int((float(SLIDE_H) - render_h) / 2)

        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
        img.save(tmp.name, format="JPEG", quality=78, optimize=True)

    slide.shapes.add_picture(str(tmp.name), left, top, width=render_w, height=render_h)
    import os; os.unlink(tmp.name)


def add_full_bleed_image(slide, image_path: Optional[Path]) -> None:
    if not image_path or not image_path.exists():
        return
    SLIDE_W_PX, SLIDE_H_PX = 1280, 720
    try:
        with Image.open(image_path) as im:
            img = im.convert("RGB")
            img_ratio = img.width / img.height
            slide_ratio = SLIDE_W_PX / SLIDE_H_PX
            if img_ratio > slide_ratio:
                new_h = SLIDE_H_PX
                new_w = int(new_h * img_ratio)
            else:
                new_w = SLIDE_W_PX
                new_h = int(new_w / img_ratio)
            img = img.resize((new_w, new_h), Image.LANCZOS)
            lc = (new_w - SLIDE_W_PX) // 2
            if img_ratio <= slide_ratio:
                # Portrait/square image: bias crop toward top so faces stay in frame
                tc = int((new_h - SLIDE_H_PX) * 0.25)
            else:
                tc = (new_h - SLIDE_H_PX) // 2
            img = img.crop((lc, tc, lc + SLIDE_W_PX, tc + SLIDE_H_PX))
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
            img.save(tmp.name, format="JPEG", quality=88, optimize=True)
        slide.shapes.add_picture(tmp.name, 0, 0, width=SLIDE_W, height=SLIDE_H)
        os.unlink(tmp.name)
    except Exception as e:
        print(f"⚠️ Full bleed image error: {e}")


def add_title_poster_image(slide, image_path: Optional[Path]) -> None:
    add_blur_background(slide, image_path)
    add_center_image(slide, image_path, scale_factor=0.82)


def add_top_rule(slide) -> None:
    return


def add_title_text(slide, text: str) -> None:
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(0.58), Inches(12.0), Inches(0.45))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean(text)
    run.font.name = _theme_font()
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = rgb(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def add_text_box(slide, left, top, width, height, text: str, *, font_size: int = 18,
                 align=PP_ALIGN.LEFT, fill_transparency: float = 0.22) -> None:
    accent = _active_theme["accent"]
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = rgb(14, 14, 18)
    box.fill.transparency = fill_transparency
    box.line.color.rgb = rgb(*accent)
    box.line.width = Pt(1.2)

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.24)
    tf.margin_right = Inches(0.24)
    tf.margin_top = Inches(0.16)
    tf.margin_bottom = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean(text)
    run.font.name = _theme_font()
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = rgb(255, 255, 255)
    p.alignment = align


def add_panel_text(slide, left, top_title, panel_w, slide_title: str, body: str, font_size: int = 17) -> None:
    """Body text directly on a dark panel — no box, no border. Title + thin rule + body."""
    accent = _active_theme["accent"]
    tx_w = panel_w - int(Inches(0.56))
    x = left + int(Inches(0.28))

    # Title
    tx_title = slide.shapes.add_textbox(x, top_title, tx_w, Inches(0.52))
    tf = tx_title.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; run = p.add_run()
    run.text = clean(slide_title.split("(")[0].strip())
    run.font.name = _theme_font(); run.font.size = Pt(15)
    run.font.bold = True; run.font.color.rgb = rgb(*accent)
    p.alignment = PP_ALIGN.LEFT

    # Thin accent rule under title
    rule_y = top_title + int(Inches(0.6))
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, rule_y, tx_w, int(Inches(0.025)))
    rule.fill.solid(); rule.fill.fore_color.rgb = rgb(*accent)
    rule.fill.transparency = 0.5; rule.line.fill.background()

    # Body text rectangle — vertically centered in remaining space
    body_y = rule_y + int(Inches(0.12))
    body_h = SLIDE_H - body_y - int(Inches(0.36))
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, body_y, tx_w, body_h)
    rect.fill.background(); rect.line.fill.background()
    tf2 = rect.text_frame; tf2.clear(); tf2.word_wrap = True
    tf2.margin_left = Inches(0.0); tf2.margin_right = Inches(0.1)
    tf2.margin_top = Inches(0.18); tf2.margin_bottom = Inches(0.0)
    tf2.vertical_anchor = MSO_ANCHOR.TOP
    p2 = tf2.paragraphs[0]; run2 = p2.add_run()
    run2.text = clean(body)
    run2.font.name = _theme_font(); run2.font.size = Pt(font_size)
    run2.font.bold = True; run2.font.color.rgb = rgb(240, 240, 240)
    p2.alignment = PP_ALIGN.LEFT
    p2.line_spacing = Pt(font_size * 1.45)


def add_cinematic_caption(slide, body: str, font_size: int = 18) -> None:
    """Full-width dark band anchored at the bottom — no border, text sits on the image."""
    if not body:
        return
    band_h = Inches(1.55)
    band_top = SLIDE_H - band_h
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, band_top, SLIDE_W, band_h)
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(6, 6, 8)
    band.fill.transparency = 0.18
    band.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.6), band_top, SLIDE_W - Inches(1.2), band_h)
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.18)
    tf.margin_bottom = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean(body)
    run.font.name = _theme_font()
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = rgb(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def _auto_font_size(text: str, base: int) -> int:
    n = len(text)
    if n > 300: return max(10, base - 5)
    if n > 220: return max(11, base - 4)
    if n > 160: return max(12, base - 3)
    if n > 110: return max(13, base - 2)
    if n > 70:  return max(14, base - 1)
    return base


def build_slide_split_panel(slide, image_path: Optional[Path], slide_title: str, body: str) -> None:
    """Layout B — image fills left 55%, dark text panel right 45%."""
    add_base_background(slide)
    accent = _active_theme["accent"]

    panel_w = int(float(SLIDE_W) * 0.55)
    panel_h = int(float(SLIDE_H))

    # Pixel dimensions for image processing (EMU units above are not pixels)
    panel_w_px = 704
    panel_h_px = 720

    if image_path and image_path.exists():
        try:
            with Image.open(image_path) as im:
                img = im.convert("RGB")
                img.thumbnail((panel_w_px * 2, panel_h_px * 2))
                img_ratio = img.width / img.height
                panel_ratio = panel_w_px / panel_h_px
                if img_ratio > panel_ratio:
                    new_h = panel_h_px
                    new_w = int(new_h * img_ratio)
                else:
                    new_w = panel_w_px
                    new_h = int(new_w / img_ratio)
                img = img.resize((new_w, new_h), Image.LANCZOS)
                lc = (new_w - panel_w_px) // 2
                tc = (new_h - panel_h_px) // 2
                img = img.crop((lc, tc, lc + panel_w_px, tc + panel_h_px))
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
                img.save(tmp.name, format="JPEG", quality=82, optimize=True)
            slide.shapes.add_picture(str(tmp.name), 0, 0, width=panel_w, height=SLIDE_H)
            os.unlink(tmp.name)
        except Exception:
            pass

    # Dark right panel
    right_x = int(float(SLIDE_W) * 0.54)
    right_w = int(float(SLIDE_W) - right_x)
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, right_x, 0, right_w, SLIDE_H)
    panel.fill.solid()
    panel.fill.fore_color.rgb = rgb(10, 10, 14)
    panel.fill.transparency = 0.0
    panel.line.fill.background()

    # Accent divider line
    div = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, right_x, 0, Inches(0.04), SLIDE_H)
    div.fill.solid()
    div.fill.fore_color.rgb = rgb(*accent)
    div.fill.transparency = 0.35
    div.line.fill.background()

    # Title + rule + body directly on dark panel (no box)
    font_size = _auto_font_size(body, base=17)
    add_panel_text(slide, right_x, Inches(0.52), right_w, slide_title, body, font_size=font_size)


def build_slide_person_attached(slide, image_path: Optional[Path], name: str, role: str, credits_line: str) -> None:
    """Person Attached slide — portrait photo left, role/name/credits right."""
    add_base_background(slide)
    accent = _active_theme["accent"]

    panel_w = int(float(SLIDE_W) * 0.50)
    panel_w_px, panel_h_px = 640, 720

    if image_path and image_path.exists():
        try:
            with Image.open(image_path) as im:
                img = im.convert("RGB")
                img_ratio = img.width / img.height
                panel_ratio = panel_w_px / panel_h_px
                if img_ratio > panel_ratio:
                    new_h = panel_h_px
                    new_w = int(new_h * img_ratio)
                else:
                    new_w = panel_w_px
                    new_h = int(new_w / img_ratio)
                img = img.resize((new_w, new_h), Image.LANCZOS)
                lc = (new_w - panel_w_px) // 2
                tc = max(0, min(int(new_h * 0.05), new_h - panel_h_px))
                img = img.crop((lc, tc, lc + panel_w_px, tc + panel_h_px))
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
                img.save(tmp.name, format="JPEG", quality=85, optimize=True)
            slide.shapes.add_picture(str(tmp.name), 0, 0, width=panel_w, height=SLIDE_H)
            os.unlink(tmp.name)
        except Exception as e:
            print(f"⚠️ Person photo error: {e}")

    right_x = int(float(SLIDE_W) * 0.49)
    right_w = int(float(SLIDE_W) - right_x)

    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, right_x, 0, right_w, SLIDE_H)
    panel.fill.solid(); panel.fill.fore_color.rgb = rgb(10, 10, 14)
    panel.fill.transparency = 0.0; panel.line.fill.background()

    div = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, right_x, 0, Inches(0.04), SLIDE_H)
    div.fill.solid(); div.fill.fore_color.rgb = rgb(*accent)
    div.fill.transparency = 0.35; div.line.fill.background()

    pad_x = right_x + int(Inches(0.32))
    text_w = right_w - int(Inches(0.64))

    # Role label
    tb_role = slide.shapes.add_textbox(pad_x, Inches(1.1), text_w, Inches(0.42))
    tf = tb_role.text_frame; tf.clear()
    p = tf.paragraphs[0]; run = p.add_run()
    run.text = clean(role).upper()
    run.font.name = _theme_font(); run.font.size = Pt(13)
    run.font.bold = True; run.font.color.rgb = rgb(*accent)

    # Name
    tb_name = slide.shapes.add_textbox(pad_x, Inches(1.62), text_w, Inches(1.4))
    tf2 = tb_name.text_frame; tf2.clear(); tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; run2 = p2.add_run()
    run2.text = clean(name)
    run2.font.name = _theme_font(); run2.font.size = Pt(30)
    run2.font.bold = True; run2.font.color.rgb = rgb(255, 255, 255)

    # Credits line in rounded card
    if credits_line:
        card_y = Inches(3.2)
        card_h = SLIDE_H - card_y - int(Inches(0.55))
        add_text_box(slide, pad_x, card_y, text_w, card_h, credits_line, font_size=15, align=PP_ALIGN.LEFT, fill_transparency=0.22)


def build_slide_text_only(slide, slide_title: str, body: str) -> None:
    """Text-only layout — no image. Big centered body text fills the slide."""
    add_base_background(slide)
    accent = _active_theme["accent"]
    add_top_rule(slide)

    # Title band
    tx = slide.shapes.add_textbox(Inches(1.0), Inches(0.52), Inches(11.3), Inches(0.7))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean(slide_title.split("(")[0].strip())
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = rgb(*accent)
    p.alignment = PP_ALIGN.CENTER

    # Large body — fills most of slide
    font_size = _auto_font_size(body, base=22)
    font_size = max(font_size, 20)
    add_text_box(slide, Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.4),
                 body, font_size=font_size, align=PP_ALIGN.CENTER, fill_transparency=0.0)


def build_slide_editorial(slide, image_path: Optional[Path], slide_title: str, body: str) -> None:
    """Layout C — image floats center-top, title below it, wide body box at bottom."""
    add_base_background(slide)
    accent = _active_theme["accent"]

    if image_path and image_path.exists():
        try:
            with Image.open(image_path) as _probe:
                _is_portrait = _probe.height > _probe.width
        except Exception:
            _is_portrait = False
        if _is_portrait:
            add_blur_background(slide, image_path)
            add_center_image(slide, image_path, scale_factor=0.72)
        else:
            add_center_image(slide, image_path, scale_factor=0.52)

    # Title centered below image area
    tx = slide.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(11.3), Inches(0.6))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean(slide_title.split("(")[0].strip())
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = rgb(*accent)
    p.alignment = PP_ALIGN.CENTER

    # Wide body box at bottom
    font_size = _auto_font_size(body, base=16)
    add_text_box(slide, Inches(0.7), Inches(5.1), Inches(11.93), Inches(2.0),
                 body, font_size=font_size, align=PP_ALIGN.CENTER, fill_transparency=0.18)


def build_slide_split_right(slide, image_path: Optional[Path], slide_title: str, body: str) -> None:
    """Layout D — dark text panel left 45%, image fills right 55%."""
    add_base_background(slide)
    accent = _active_theme["accent"]

    panel_w = int(float(SLIDE_W) * 0.45)
    img_x = int(float(SLIDE_W) * 0.45)
    img_w = int(float(SLIDE_W) - img_x)
    panel_w_px, panel_h_px = 576, 720

    if image_path and image_path.exists():
        try:
            with Image.open(image_path) as im:
                img = im.convert("RGB")
                img_ratio = img.width / img.height
                panel_ratio = img_w / SLIDE_H
                if img_ratio > panel_ratio:
                    new_h = int(SLIDE_H / 914400 * 96)
                    new_w = int(new_h * img_ratio)
                else:
                    new_w = panel_w_px
                    new_h = int(new_w / img_ratio)
                new_w = max(new_w, panel_w_px)
                new_h = max(new_h, panel_h_px)
                img = img.resize((new_w, new_h), Image.LANCZOS)
                lc = (new_w - panel_w_px) // 2
                tc = (new_h - panel_h_px) // 2
                img = img.crop((lc, tc, lc + panel_w_px, tc + panel_h_px))
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
                img.save(tmp.name, format="JPEG", quality=82, optimize=True)
            slide.shapes.add_picture(str(tmp.name), img_x, 0, width=img_w, height=SLIDE_H)
            os.unlink(tmp.name)
        except Exception:
            pass

    # Dark left panel
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, panel_w, SLIDE_H)
    panel.fill.solid()
    panel.fill.fore_color.rgb = rgb(10, 10, 14)
    panel.fill.transparency = 0.0
    panel.line.fill.background()

    # Accent divider line (right edge of panel)
    div = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, panel_w - int(Inches(0.04)), 0, int(Inches(0.04)), SLIDE_H)
    div.fill.solid()
    div.fill.fore_color.rgb = rgb(*accent)
    div.fill.transparency = 0.35
    div.line.fill.background()

    # Title + rule + body directly on dark panel (no box)
    font_size = _auto_font_size(body, base=17)
    add_panel_text(slide, 0, Inches(0.52), panel_w, slide_title, body, font_size=font_size)


def build_slide_quote_overlay(slide, image_path: Optional[Path], slide_title: str, body: str) -> None:
    """Layout E — full bleed image, small title top-left, large body floats center-lower."""
    add_base_background(slide)
    add_full_bleed_image(slide, image_path)
    accent = _active_theme["accent"]

    # Small title top-left
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(0.38), Inches(10.0), Inches(0.52))
    tf = tx.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; run = p.add_run()
    run.text = clean(slide_title.split("(")[0].strip().upper())
    run.font.name = _theme_font(); run.font.size = Pt(12)
    run.font.bold = True; run.font.color.rgb = rgb(*accent)
    p.alignment = PP_ALIGN.LEFT

    if body:
        fs = _auto_font_size(body, base=22)
        fs = max(fs, 17)
        add_text_box(slide, Inches(0.8), Inches(2.6), Inches(11.73), Inches(3.6),
                     body, font_size=fs, align=PP_ALIGN.CENTER, fill_transparency=0.26)


def build_slide_bottom_card(slide, image_path: Optional[Path], slide_title: str, body: str) -> None:
    """Layout F — full bleed image, prominent bottom card with accent rule, title + body."""
    add_base_background(slide)
    add_full_bleed_image(slide, image_path)
    accent = _active_theme["accent"]

    card_h = Inches(2.75)
    card_top = SLIDE_H - card_h

    # Card background
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, card_top, SLIDE_W, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = rgb(6, 6, 10)
    card.fill.transparency = 0.10
    card.line.fill.background()

    # Accent rule at top of card
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, card_top, SLIDE_W, int(Inches(0.05)))
    rule.fill.solid()
    rule.fill.fore_color.rgb = rgb(*accent)
    rule.fill.transparency = 0.0
    rule.line.fill.background()

    # Title inside card
    tx_title = slide.shapes.add_textbox(Inches(0.6), card_top + int(Inches(0.18)), Inches(12.13), Inches(0.56))
    tf = tx_title.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; run = p.add_run()
    run.text = clean(slide_title.split("(")[0].strip())
    run.font.name = _theme_font(); run.font.size = Pt(15)
    run.font.bold = True; run.font.color.rgb = rgb(*accent)
    p.alignment = PP_ALIGN.LEFT

    if body:
        fs = _auto_font_size(body, base=16)
        tx_body = slide.shapes.add_textbox(Inches(0.6), card_top + int(Inches(0.86)), Inches(12.13), Inches(1.72))
        tf2 = tx_body.text_frame; tf2.clear(); tf2.word_wrap = True
        tf2.margin_left = Inches(0.08)
        p2 = tf2.paragraphs[0]; run2 = p2.add_run()
        run2.text = clean(body)
        run2.font.name = _theme_font(); run2.font.size = Pt(fs)
        run2.font.bold = True; run2.font.color.rgb = rgb(238, 238, 242)
        p2.alignment = PP_ALIGN.LEFT


def place_text_by_stage(slide, stage: str, layout: str, body: str) -> None:
    stage = clean(stage).lower()
    layout = clean(layout).lower()
    if not body:
        return
    fs = _auto_font_size(body, 18)
    if layout in {"title", "closing"}:
        fs = _auto_font_size(body, 19)
    add_cinematic_caption(slide, body, font_size=fs)


def _prefetch_slide_image(args: tuple) -> tuple:
    """Run in a thread pool — resolves/generates the image for one slide."""
    idx, slide_info, visuals_dir, deck_title, brain_output = args
    slide_title = clean(slide_info.get("title"))
    body = clean(slide_info.get("body"))
    layout = clean(slide_info.get("layout", ""))
    slide_number = int(slide_info.get("slide_number", idx))

    explicit_path_str = str(slide_info.get("image_path") or "").strip()
    image_source_hint = str(slide_info.get("image_source") or "").strip()

    if explicit_path_str == "__none__" or image_source_hint == "text_only":
        return idx, None, "text_only"

    # Person attached: download headshot from TMDb URL
    person_photo_url = str(slide_info.get("person_photo_url") or "").strip()
    if person_photo_url and layout.lower() == "person_attached":
        try:
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
            urllib.request.urlretrieve(person_photo_url, tmp.name)
            return idx, Path(tmp.name), "tmdb"
        except Exception as e:
            print(f"⚠️ TMDb photo download failed: {e}")
            return idx, None, "tmdb_failed"

    if explicit_path_str:
        explicit = Path(explicit_path_str)
        if not explicit.is_absolute():
            explicit = (APP_DIR / explicit).resolve()
        if explicit.exists():
            return idx, explicit, str(slide_info.get("image_source") or "user_selected")

    img, src = find_image_for_slide(
        visuals_dir=visuals_dir,
        deck_title=deck_title,
        slide_title=slide_title if layout.lower() != "title" else deck_title,
        slide_number=slide_number,
        brain_output=brain_output,
        last_used_name="",
        slide_body=body,
    )
    return idx, img, src


def _fetch_tmdb_poster(title: str, cache_dir: Path) -> Optional[Path]:
    """Download a movie poster from TMDb. Returns local path or None."""
    if not TMDB_API_KEY:
        return None
    import urllib.parse
    safe = re.sub(r"[^a-z0-9_]", "_", title.lower())[:40]
    cache_path = cache_dir / f"tmdb_{safe}.jpg"
    if cache_path.exists():
        return cache_path
    try:
        search_url = (
            f"https://api.themoviedb.org/3/search/movie"
            f"?query={urllib.parse.quote(title)}&include_adult=false"
        )
        req = urllib.request.Request(
            search_url,
            headers={"Authorization": f"Bearer {TMDB_API_KEY}", "Accept": "application/json"},
        )
        with urllib.request.urlopen(req, timeout=10) as resp:
            data = json.loads(resp.read())
        results = data.get("results", [])
        if not results:
            return None
        poster_path = results[0].get("poster_path")
        if not poster_path:
            return None
        cache_path.parent.mkdir(parents=True, exist_ok=True)
        urllib.request.urlretrieve(f"{TMDB_IMAGE_BASE}{poster_path}", cache_path)
        return cache_path
    except Exception as e:
        print(f"⚠️ TMDb fetch failed for '{title}': {e}")
        return None


def _build_comp_poster_strip_slide(slide, poster_paths: list, comp_titles: list) -> None:
    """Film strip layout — real movie posters side by side with gold frame borders."""
    add_base_background(slide)
    accent = _active_theme["accent"]

    label_tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(6.0), Inches(0.45))
    tf = label_tx.text_frame; tf.clear()
    p = tf.paragraphs[0]; run = p.add_run()
    run.text = "COMPARABLES"
    run.font.name = _theme_font(); run.font.size = Pt(13)
    run.font.bold = True; run.font.color.rgb = rgb(*accent)

    n = max(len(poster_paths), 1)
    margin_x  = Inches(0.45)
    margin_top = Inches(0.88)
    margin_bot = Inches(0.78)
    spacing    = Inches(0.22)
    total_w    = SLIDE_W - 2 * margin_x - (n - 1) * spacing
    poster_w   = int(total_w / n)
    poster_h   = int(SLIDE_H - margin_top - margin_bot)

    for i, (ppath, ptitle) in enumerate(zip(poster_paths, comp_titles)):
        x = int(margin_x + i * (poster_w + spacing))
        y = int(margin_top)

        border = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            x - int(Inches(0.04)), y - int(Inches(0.04)),
            poster_w + int(Inches(0.08)), poster_h + int(Inches(0.08)),
        )
        border.fill.solid(); border.fill.fore_color.rgb = rgb(*accent)
        border.fill.transparency = 0.55; border.line.fill.background()

        if ppath and ppath.exists():
            try:
                PW = int(poster_w / 914400 * 96)
                PH = int(poster_h / 914400 * 96)
                with Image.open(ppath) as im:
                    img = im.convert("RGB")
                    ir = img.width / img.height
                    pr = PW / PH
                    if ir > pr:
                        nh = PH; nw = int(nh * ir)
                    else:
                        nw = PW; nh = int(nw / ir)
                    img = img.resize((nw, nh), Image.LANCZOS)
                    lc = (nw - PW) // 2; tc = (nh - PH) // 2
                    img = img.crop((lc, tc, lc + PW, tc + PH))
                    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
                    img.save(tmp.name, format="JPEG", quality=82, optimize=True)
                slide.shapes.add_picture(tmp.name, x, y, width=poster_w, height=poster_h)
                os.unlink(tmp.name)
            except Exception as e:
                print(f"⚠️ Comp poster render error: {e}")

        if ptitle:
            tb = slide.shapes.add_textbox(x, y + poster_h + int(Inches(0.08)), poster_w, int(Inches(0.5)))
            tf2 = tb.text_frame; tf2.clear(); tf2.word_wrap = True
            p2 = tf2.paragraphs[0]; run2 = p2.add_run()
            run2.text = clean(ptitle)
            run2.font.name = _theme_font(); run2.font.size = Pt(11)
            run2.font.bold = True; run2.font.color.rgb = rgb(220, 220, 220)
            p2.alignment = PP_ALIGN.CENTER


def _build_poster_cover_slide(slide, image_path: Optional[Path], title: str, tagline: str, brain_output: dict) -> None:
    """Movie poster style cover — full-bleed image, dark gradient, large gold title, italic tagline."""
    accent = _active_theme["accent"]

    if image_path and image_path.exists():
        add_full_bleed_image(slide, image_path)
    else:
        add_base_background(slide)

    # Dark gradient overlay covering bottom 55%
    overlay_h = int(float(SLIDE_H) * 0.55)
    overlay_y = int(float(SLIDE_H) - overlay_h)
    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, overlay_y, SLIDE_W, overlay_h)
    overlay.fill.solid(); overlay.fill.fore_color.rgb = rgb(4, 4, 6)
    overlay.fill.transparency = 0.05; overlay.line.fill.background()

    # Thin gold rule
    rule_y = int(float(SLIDE_H) * 0.52)
    rule = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, int(Inches(0.7)), rule_y,
        int(SLIDE_W - Inches(1.4)), int(Inches(0.04)),
    )
    rule.fill.solid(); rule.fill.fore_color.rgb = rgb(*accent)
    rule.fill.transparency = 0.0; rule.line.fill.background()

    # Large title
    title_clean = clean(title).upper()
    title_fs = 52 if len(title_clean) <= 18 else (42 if len(title_clean) <= 28 else 34)
    title_y = int(float(SLIDE_H) * 0.54)
    tx_title = slide.shapes.add_textbox(int(Inches(0.5)), title_y, int(SLIDE_W - Inches(1.0)), int(Inches(1.5)))
    tf = tx_title.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; run = p.add_run()
    run.text = title_clean
    run.font.name = _theme_font(); run.font.size = Pt(title_fs)
    run.font.bold = True; run.font.color.rgb = rgb(*accent)
    p.alignment = PP_ALIGN.CENTER

    # Italic tagline
    if tagline:
        tag_y = title_y + int(Inches(1.5))
        tx_tag = slide.shapes.add_textbox(int(Inches(1.0)), tag_y, int(SLIDE_W - Inches(2.0)), int(Inches(0.65)))
        tf2 = tx_tag.text_frame; tf2.clear(); tf2.word_wrap = True
        p2 = tf2.paragraphs[0]; run2 = p2.add_run()
        run2.text = clean(tagline)
        run2.font.name = _theme_font(); run2.font.size = Pt(17)
        run2.font.bold = False; run2.font.italic = True
        run2.font.color.rgb = rgb(210, 210, 210)
        p2.alignment = PP_ALIGN.CENTER

    # Writer credit top-left (if available)
    writer = clean(brain_output.get("writer", "") or brain_output.get("author", ""))
    if writer:
        tx_w = slide.shapes.add_textbox(int(Inches(0.55)), int(Inches(0.28)), int(Inches(5.0)), int(Inches(0.42)))
        tf3 = tx_w.text_frame; tf3.clear()
        p3 = tf3.paragraphs[0]; run3 = p3.add_run()
        run3.text = f"Written by {writer}"
        run3.font.name = _theme_font(); run3.font.size = Pt(11)
        run3.font.color.rgb = rgb(180, 180, 180)


def build_presentation(
    slide_plan_path: Path,
    visuals_dir: Path,
    output_dir: Path,
    label: str = "",
    uid: str = ""
) -> Path:
    global _active_theme
    reset_image_selection_state()
    plan = load_json(slide_plan_path)
    brain_output = load_brain_output(output_dir)

    layout_strategy = brain_output.get("layout_strategy") or {}
    layout_style = (layout_strategy.get("layout_style") or "cinematic_grounded").strip()
    composition_bias = (layout_strategy.get("composition_bias") or "image_forward").strip()
    _active_theme = LAYOUT_THEMES.get(layout_style, LAYOUT_THEMES["cinematic_grounded"])
    print(f"🎨 Layout theme: {layout_style} | Composition: {composition_bias}")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    deck_title = clean(plan.get("title", "Project"))
    manifest: list[dict] = []

    # Pre-generate all FAL images in parallel before building slides
    slides_list = plan.get("slides", [])
    prefetch_args = [
        (idx, slide_info, visuals_dir, deck_title, brain_output)
        for idx, slide_info in enumerate(slides_list, start=1)
    ]
    prefetched: dict[int, tuple] = {}
    print(f"🖼️  Pre-fetching images for {len(prefetch_args)} slides in parallel...")
    with concurrent.futures.ThreadPoolExecutor(max_workers=8) as pool:
        for idx, img, src in pool.map(_prefetch_slide_image, prefetch_args):
            prefetched[idx] = (img, src)
    print("✅ Image pre-fetch complete")

    last_used_image_name = ""

    for idx, slide_info in enumerate(slides_list, start=1):
        slide_title = clean(slide_info.get("title"))
        body = clean(slide_info.get("body"))
        layout = clean(slide_info.get("layout"))
        stage = clean(slide_info.get("stage"))
        slide_number = int(slide_info.get("slide_number", idx))

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        image_for_slide, image_source = prefetched.get(idx, (None, "none"))

        if image_for_slide:
            last_used_image_name = image_for_slide.name
            _mark_image_used(image_for_slide)

        stage_lower = clean(stage).lower()
        layout_lower = clean(layout).lower()
        _stitle = slide_title.split("(")[0].strip()

        if layout_lower == "person_attached":
            build_slide_person_attached(
                slide,
                image_for_slide,
                clean(slide_info.get("person_name", "") or _stitle),
                clean(slide_info.get("person_role", "") or "Attached"),
                clean(slide_info.get("person_credits_line", "") or body),
            )

        elif image_source == "text_only":
            build_slide_text_only(slide, slide_title, body)

        elif layout_lower == "title":
            if image_source not in {"poster", ""} and image_for_slide is not None:
                _title_img = image_for_slide
            else:
                _title_img = Path(POSTER_PATH) if POSTER_PATH else image_for_slide
            _tagline = clean(brain_output.get("tagline", "") or "")
            if not _tagline:
                _logline = clean(brain_output.get("logline", "") or "")
                _tagline = _logline[:90] if _logline else ""
            _build_poster_cover_slide(slide, _title_img, deck_title, _tagline, brain_output)

        elif layout_lower in {"character_focus", "split_left_text"} or stage_lower == "character":
            build_slide_split_panel(slide, image_for_slide, _stitle, body)

        elif layout_lower == "split_right_text":
            build_slide_split_right(slide, image_for_slide, _stitle, body)

        elif layout_lower == "quote_overlay" or stage_lower in {"world", "tone", "themes"}:
            build_slide_quote_overlay(slide, image_for_slide, _stitle, body)

        elif layout_lower == "bottom_story_card" or stage_lower in {"engine", "setup", "aftermath", "why_now"}:
            build_slide_bottom_card(slide, image_for_slide, _stitle, body)

        elif normalize_key(slide_title) in {"comparables", "comparable films", "comps"} and TMDB_API_KEY:
            _comp_raw = brain_output.get("comparables", [])
            if isinstance(_comp_raw, str):
                _comp_raw = [t.strip() for t in _comp_raw.split(",") if t.strip()]
            _comp_titles = [clean(str(t)) for t in _comp_raw if t][:4]
            _cache_dir = APP_DIR / "generated_images" / EVOLUM_SESSION_ID / "tmdb_posters"
            _poster_paths = [_fetch_tmdb_poster(t, _cache_dir) for t in _comp_titles]
            _build_comp_poster_strip_slide(slide, _poster_paths, _comp_titles)

        elif layout_lower == "clean_grid" or stage_lower == "market":
            build_slide_editorial(slide, image_for_slide, _stitle, body)

        elif layout_lower == "hero_full_bleed" or stage_lower == "closing":
            if stage_lower == "title":
                _title_img = image_for_slide
                if not _title_img and POSTER_PATH:
                    _title_img = Path(POSTER_PATH)
                _tagline = clean(brain_output.get("tagline", "") or "")
                if not _tagline:
                    _logline = clean(brain_output.get("logline", "") or "")
                    _tagline = _logline[:90] if _logline else ""
                _build_poster_cover_slide(slide, _title_img, deck_title, _tagline, brain_output)
            else:
                add_base_background(slide)
                add_full_bleed_image(slide, image_for_slide)
                add_title_text(slide, deck_title if stage_lower == "closing" else _stitle)
                place_text_by_stage(slide, stage, layout, body)

        else:
            # Default — let composition_bias from brain guide the choice
            if composition_bias == "split_text_image":
                build_slide_split_panel(slide, image_for_slide, _stitle, body)
            elif composition_bias == "illustrative":
                build_slide_editorial(slide, image_for_slide, _stitle, body)
            else:
                # Full bleed with bottom card (richer than plain caption)
                build_slide_bottom_card(slide, image_for_slide, _stitle, body)

        resolved_image_options = resolve_image_options_for_slide(
            visuals_dir=visuals_dir,
            slide_info=slide_info,
            image_for_slide=image_for_slide,
            image_source=image_source,
            slide_title=slide_title,
        )

        user_selected_option_id = str(slide_info.get("selected_option_id") or "").strip()
        user_image_url = str(slide_info.get("image_url") or "").strip()
        user_image_name = str(slide_info.get("image_name") or "").strip()

        manifest.append({
            "slide_number": slide_number,
            "title": slide_title,
            "subtitle": str(slide_info.get("subtitle") or "").strip(),
            "body": body,
            "layout": layout,
            "stage": stage,
            "image_path": "__none__" if image_source == "text_only" else (str(image_for_slide) if image_for_slide else ""),
            "image_name": user_image_name or (image_for_slide.name if image_for_slide else ""),
            "image_url": user_image_url,
            "image_source": image_source,
            "image_query": slide_info.get("image_query", ""),
            "image_tags": slide_info.get("image_tags", []),
            "image_score": slide_info.get("image_score", 0),
            "image_options": resolved_image_options,
            "selected_option_id": user_selected_option_id or (resolved_image_options[0].get("option_id", "selected") if resolved_image_options else ""),
        })

    out_path = next_output_path(output_dir, label=label, uid=uid)
    prs.save(str(out_path))
    prefix = f"{uid}_" if uid else ""
    manifest_name = f"{prefix}latest_deck_manifest_{label}.json" if label else f"{prefix}latest_deck_manifest.json"
    manifest_path = output_dir / manifest_name
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    print(f"📦 Deck manifest created: {manifest_path}")
    print(f"✅ Pitch deck created: {out_path}")
    _write_deck_builder_tokens()
    return out_path


def parse_args():
    parser = argparse.ArgumentParser()
    parser.add_argument("slide_plan", nargs="?", help="Path to slide_plan.json")
    parser.add_argument("--project", help="Project/app directory containing slide_plan.json and visuals/")
    parser.add_argument("--label", default="", help="Output label (e.g. 'producer') for file naming")
    parser.add_argument("--uid", default="", help="User ID for per-user manifest file naming")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    slide_plan_path, visuals_dir, output_dir = resolve_paths(args)

    if not slide_plan_path.exists():
        print(f"Slide plan not found: {slide_plan_path}")
        raise SystemExit(1)

    build_presentation(slide_plan_path, visuals_dir, output_dir, label=args.label, uid=args.uid)


if __name__ == "__main__":
    main()
